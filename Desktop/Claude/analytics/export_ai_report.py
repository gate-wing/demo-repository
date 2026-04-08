"""
AI参照トラフィック横断レポート生成スクリプト

ai_traffic_data.json を読み込み、以下を出力します：
  - ai_traffic_report_YYYYMM.md    : Markdownレポート
  - ai_traffic_report_YYYYMM.pptx  : PowerPointレポート（グラフ付き）

実行方法:
  python -X utf8 export_ai_report.py
"""

import sys, os, json
if sys.platform == "win32":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib import rcParams
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tempfile

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ── カラー設定 ─────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x1A, 0x23, 0x3A)
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY   = RGBColor(0x88, 0x88, 0x88)

ENGINE_COLORS = {
    "ChatGPT":           "#10a37f",
    "Gemini":            "#4285F4",
    "Perplexity":        "#20B2AA",
    "Microsoft Copilot": "#00B4EF",
    "Claude":            "#D97706",
}
DEFAULT_COLOR = "#999999"

# ── フォント設定（Windows対応） ────────────────────────────────
for font in ["Meiryo", "Yu Gothic", "MS Gothic", "IPAexGothic"]:
    try:
        rcParams["font.family"] = font
        break
    except Exception:
        pass
rcParams["axes.unicode_minus"] = False


# ── データ読み込み ─────────────────────────────────────────────
def load_data():
    path = os.path.join(BASE_DIR, "ai_traffic_data.json")
    if not os.path.exists(path):
        print("❌ ai_traffic_data.json が見つかりません。先に fetch_ai_traffic.py を実行してください。")
        sys.exit(1)
    with open(path, encoding="utf-8") as f:
        return json.load(f)


# ── 集計ヘルパー ───────────────────────────────────────────────
def total_by_engine(data):
    """全プロパティ・全期間のAIエンジン別合計を返す"""
    totals = {}
    for prop in data["properties"]:
        for month in prop["monthly_data"]:
            for engine, vals in month["by_engine"].items():
                totals[engine] = totals.get(engine, 0) + vals["sessions"]
    return totals


def monthly_total_all(data):
    """全プロパティ合算の月別合計 {month: total_sessions}"""
    monthly = {}
    for prop in data["properties"]:
        for month in prop["monthly_data"]:
            lbl = month["month"]
            monthly[lbl] = monthly.get(lbl, 0) + month["total_sessions"]
    return dict(sorted(monthly.items()))


def monthly_by_engine_all(data):
    """全プロパティ合算の月別・エンジン別 {engine: {month: sessions}}"""
    result = {}
    for prop in data["properties"]:
        for month in prop["monthly_data"]:
            lbl = month["month"]
            for engine, vals in month["by_engine"].items():
                if engine not in result:
                    result[engine] = {}
                result[engine][lbl] = result[engine].get(lbl, 0) + vals["sessions"]
    return result


# ── グラフ生成 ─────────────────────────────────────────────────
def make_chart_monthly_total(data, months):
    """全体の月別AI参照合計グラフ"""
    monthly = monthly_total_all(data)
    labels  = list(monthly.keys())
    values  = list(monthly.values())

    fig, ax = plt.subplots(figsize=(10, 4))
    bars = ax.bar(labels, values, color="#FF6B35", alpha=0.85, width=0.6)
    ax.set_title("AI参照セッション数 推移（全店舗合計）", fontsize=13, pad=10)
    ax.set_xlabel("月", fontsize=10)
    ax.set_ylabel("セッション数", fontsize=10)
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    plt.xticks(rotation=45, ha="right", fontsize=8)
    plt.tight_layout()

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def make_chart_engine_pie(data):
    """AIエンジン別シェア 円グラフ"""
    totals = total_by_engine(data)
    # 0件のエンジンは除外
    totals = {k: v for k, v in totals.items() if v > 0}
    if not totals:
        return None

    labels = list(totals.keys())
    values = list(totals.values())
    colors = [ENGINE_COLORS.get(l, DEFAULT_COLOR) for l in labels]

    fig, ax = plt.subplots(figsize=(6, 5))
    wedges, texts, autotexts = ax.pie(
        values, labels=labels, colors=colors,
        autopct="%1.1f%%", startangle=90,
        textprops={"fontsize": 9}
    )
    ax.set_title("AIエンジン別シェア（全店舗・全期間）", fontsize=12, pad=10)
    plt.tight_layout()

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def make_chart_engine_trend(data):
    """AIエンジン別 月別トレンド 折れ線グラフ"""
    by_engine = monthly_by_engine_all(data)
    months = sorted(next(iter(by_engine.values())).keys()) if by_engine else []

    # セッション0件のエンジンを除外
    by_engine = {e: v for e, v in by_engine.items() if sum(v.values()) > 0}
    if not by_engine:
        return None

    fig, ax = plt.subplots(figsize=(10, 4.5))
    for engine, monthly in by_engine.items():
        values = [monthly.get(m, 0) for m in months]
        color  = ENGINE_COLORS.get(engine, DEFAULT_COLOR)
        ax.plot(months, values, marker="o", label=engine, color=color, linewidth=2, markersize=5)

    ax.set_title("AIエンジン別 月別セッション推移", fontsize=13, pad=10)
    ax.set_xlabel("月", fontsize=10)
    ax.set_ylabel("セッション数", fontsize=10)
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax.legend(loc="upper left", fontsize=9)
    plt.xticks(rotation=45, ha="right", fontsize=8)
    plt.tight_layout()

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def make_chart_property_comparison(data):
    """プロパティ別 AI参照合計 横棒グラフ"""
    prop_totals = {}
    for prop in data["properties"]:
        total = sum(m["total_sessions"] for m in prop["monthly_data"])
        prop_totals[prop["display_name"]] = total

    if not any(prop_totals.values()):
        return None

    # 降順ソート
    sorted_props = sorted(prop_totals.items(), key=lambda x: x[1], reverse=True)
    names  = [p[0] for p in sorted_props]
    values = [p[1] for p in sorted_props]

    fig, ax = plt.subplots(figsize=(9, max(3, len(names) * 0.5 + 1.5)))
    bars = ax.barh(names, values, color="#1A233A", alpha=0.8)
    ax.set_title("店舗別 AI参照セッション合計（全期間）", fontsize=12, pad=10)
    ax.set_xlabel("セッション数", fontsize=10)
    ax.xaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + max(values) * 0.01, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", fontsize=9)
    plt.tight_layout()

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def make_chart_property_monthly(prop):
    """プロパティ個別の月別AI参照推移（積み上げ棒グラフ）"""
    months  = [m["month"] for m in prop["monthly_data"]]
    engines = list(prop["monthly_data"][0]["by_engine"].keys()) if prop["monthly_data"] else []
    engines = [e for e in engines if sum(m["by_engine"][e]["sessions"] for m in prop["monthly_data"]) > 0]

    if not engines:
        return None

    fig, ax = plt.subplots(figsize=(10, 4))
    bottom = [0] * len(months)
    for engine in engines:
        values = [m["by_engine"][engine]["sessions"] for m in prop["monthly_data"]]
        color  = ENGINE_COLORS.get(engine, DEFAULT_COLOR)
        ax.bar(months, values, bottom=bottom, label=engine, color=color, alpha=0.85, width=0.6)
        bottom = [b + v for b, v in zip(bottom, values)]

    ax.set_title(f"【{prop['display_name']}】AI参照月別推移", fontsize=12, pad=8)
    ax.set_xlabel("月", fontsize=10)
    ax.set_ylabel("セッション数", fontsize=10)
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax.legend(loc="upper left", fontsize=8)
    plt.xticks(rotation=45, ha="right", fontsize=8)
    plt.tight_layout()

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


# ── PPTX ヘルパー ──────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def add_image(slide, img_path, left, top, width, height):
    slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))


def add_title_bar(slide, title, subtitle=""):
    # 背景を濃紺
    set_bg(slide, DARK_NAVY)
    # タイトル
    add_text_box(slide, title, 0.4, 0.1, 9.2, 0.6,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.65, 9.2, 0.35,
                     font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    # オレンジアクセントライン
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()


# ── PPTX 生成 ─────────────────────────────────────────────────
def build_pptx(data, out_path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    fetched = data.get("fetched_at", "")[:10]
    account = data.get("target_account", "")
    engines = data.get("ai_engines", [])
    props   = data["properties"]

    prop_count = len(props)
    print(f"  対象プロパティ: {prop_count} 件")

    # --- スライド1: タイトル ─────────────────────────────────
    print("  スライド作成中: タイトル...", flush=True)
    slide = add_slide(prs)
    set_bg(slide, DARK_NAVY)
    add_text_box(slide, "AIエンジン参照トラフィック", 0.5, 1.2, 9, 0.9,
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "横断分析レポート", 0.5, 2.1, 9, 0.7,
                 font_size=22, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"{account}", 0.5, 2.85, 9, 0.5,
                 font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_text_box(slide, f"データ取得日: {fetched}", 0.5, 4.9, 9, 0.4,
                 font_size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    # --- スライド2: 全体サマリー ──────────────────────────────
    print("  スライド作成中: 全体サマリー...", flush=True)
    slide = add_slide(prs)
    set_bg(slide, LIGHT_GRAY)
    add_title_bar(slide, "全店舗サマリー", f"集計期間: 直近{data['period_months']}ヶ月")

    # KPIカード
    totals_by_engine = total_by_engine(data)
    total_all = sum(totals_by_engine.values())
    monthly   = monthly_total_all(data)
    months_list = sorted(monthly.keys())

    kpis = [
        ("AI参照\n総セッション", f"{total_all:,}", "全期間・全店舗合計"),
        ("最多月", monthly_list_max(monthly), "最もAI参照が多かった月"),
        ("対象店舗数", str(len(props)), f"プロパティ数"),
    ]
    for i, (label, value, sub) in enumerate(kpis):
        x = 0.4 + i * 3.2
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.1), Inches(2.9), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        add_text_box(slide, label, x + 0.1, 1.15, 2.7, 0.5, font_size=9, color=TEXT_GRAY)
        add_text_box(slide, value, x + 0.1, 1.55, 2.7, 0.6,
                     font_size=20, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, sub, x + 0.1, 2.35, 2.7, 0.3, font_size=8, color=TEXT_GRAY)

    # エンジン別集計テーブル
    add_text_box(slide, "■ AIエンジン別合計セッション", 0.4, 2.75, 9, 0.35,
                 font_size=10, bold=True, color=DARK_NAVY)
    col_w = 9.0 / max(len(engines), 1)
    for i, eng in enumerate(engines):
        cnt = totals_by_engine.get(eng, 0)
        x = 0.4 + i * col_w
        add_text_box(slide, eng, x, 3.1, col_w, 0.3, font_size=9, color=TEXT_GRAY)
        add_text_box(slide, f"{cnt:,}", x, 3.4, col_w, 0.35,
                     font_size=14, bold=True, color=DARK_NAVY)

    # --- スライド3: 月別推移グラフ（全体） ────────────────────
    print("  グラフ生成中: 月別推移...", flush=True)
    chart_total = make_chart_monthly_total(data, data["period_months"])
    slide = add_slide(prs)
    set_bg(slide, LIGHT_GRAY)
    add_title_bar(slide, "月別AI参照推移（全店舗合計）")
    add_image(slide, chart_total, 0.4, 1.1, 9.2, 4.2)
    os.unlink(chart_total)

    # --- スライド4: AIエンジン別シェア ───────────────────────
    print("  グラフ生成中: エンジン別シェア...", flush=True)
    chart_pie = make_chart_engine_pie(data)
    if chart_pie:
        slide = add_slide(prs)
        set_bg(slide, LIGHT_GRAY)
        add_title_bar(slide, "AIエンジン別シェア（全店舗・全期間）")
        add_image(slide, chart_pie, 1.8, 1.0, 6.4, 4.3)
        os.unlink(chart_pie)

    # --- スライド5: エンジン別トレンド ────────────────────────
    print("  グラフ生成中: エンジン別トレンド...", flush=True)
    chart_trend = make_chart_engine_trend(data)
    if chart_trend:
        slide = add_slide(prs)
        set_bg(slide, LIGHT_GRAY)
        add_title_bar(slide, "AIエンジン別 月別トレンド")
        add_image(slide, chart_trend, 0.4, 1.1, 9.2, 4.2)
        os.unlink(chart_trend)

    # --- スライド6: 店舗別比較 ────────────────────────────────
    print("  グラフ生成中: 店舗別比較...", flush=True)
    chart_comp = make_chart_property_comparison(data)
    if chart_comp:
        slide = add_slide(prs)
        set_bg(slide, LIGHT_GRAY)
        add_title_bar(slide, "店舗別 AI参照合計（全期間）")
        add_image(slide, chart_comp, 0.5, 1.1, 9.0, 4.2)
        os.unlink(chart_comp)

    # --- プロパティ個別スライド ────────────────────────────────
    for prop_idx, prop in enumerate(props, 1):
        print(f"  グラフ生成中: [{prop_idx}/{prop_count}] {prop['display_name']}...", flush=True)
        chart_prop = make_chart_property_monthly(prop)
        slide = add_slide(prs)
        set_bg(slide, LIGHT_GRAY)
        prop_total = sum(m["total_sessions"] for m in prop["monthly_data"])
        add_title_bar(slide, prop["display_name"],
                      f"AI参照合計: {prop_total:,} セッション（全期間）")
        if chart_prop:
            add_image(slide, chart_prop, 0.4, 1.1, 9.2, 4.1)
            os.unlink(chart_prop)
        else:
            add_text_box(slide, "この期間にAI参照のデータはありませんでした。",
                         0.5, 2.5, 9, 0.5, font_size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"  保存: {out_path}")


def monthly_list_max(monthly):
    """月別合計のうち最大月を返す（例: '2025-10 (42件)'）"""
    if not monthly:
        return "-"
    best = max(monthly, key=monthly.get)
    return f"{best}\n({monthly[best]:,}件)"


# ── Markdown 生成 ──────────────────────────────────────────────
def build_markdown(data, out_path):
    lines = []
    fetched = data.get("fetched_at", "")[:10]
    account = data.get("target_account", "")
    engines = data.get("ai_engines", [])
    props   = data["properties"]
    months  = data["period_months"]

    lines.append(f"# AIエンジン参照トラフィック 横断分析レポート")
    lines.append(f"")
    lines.append(f"**アカウント:** {account}  ")
    lines.append(f"**集計期間:** 直近{months}ヶ月  ")
    lines.append(f"**データ取得日:** {fetched}")
    lines.append(f"")

    # 全体サマリー
    totals_by_engine = total_by_engine(data)
    total_all = sum(totals_by_engine.values())
    lines.append(f"## 全店舗サマリー")
    lines.append(f"")
    lines.append(f"| 指標 | 値 |")
    lines.append(f"|------|-----|")
    lines.append(f"| AI参照 総セッション | {total_all:,} |")
    lines.append(f"| 対象店舗数 | {len(props)} |")
    for eng in engines:
        lines.append(f"| {eng} | {totals_by_engine.get(eng, 0):,} |")
    lines.append(f"")

    # 店舗別サマリー
    lines.append(f"## 店舗別 AI参照合計")
    lines.append(f"")
    header = "| 店舗名 | 合計 | " + " | ".join(engines) + " |"
    sep    = "|------|------|" + "------|" * len(engines)
    lines.append(header)
    lines.append(sep)
    for prop in props:
        prop_total = sum(m["total_sessions"] for m in prop["monthly_data"])
        all_engine_totals = {}
        for m in prop["monthly_data"]:
            for e, v in m["by_engine"].items():
                all_engine_totals[e] = all_engine_totals.get(e, 0) + v["sessions"]
        row = f"| {prop['display_name']} | {prop_total:,} | "
        row += " | ".join(str(all_engine_totals.get(e, 0)) for e in engines)
        row += " |"
        lines.append(row)
    lines.append(f"")

    # 月別推移
    lines.append(f"## 月別推移（全店舗合計）")
    lines.append(f"")
    monthly = monthly_total_all(data)
    lines.append(f"| 月 | AI参照セッション |")
    lines.append(f"|-----|----------------|")
    for month, total in monthly.items():
        lines.append(f"| {month} | {total:,} |")
    lines.append(f"")

    # プロパティ別詳細
    for prop in props:
        lines.append(f"## {prop['display_name']}")
        lines.append(f"")
        lines.append(f"| 月 | 合計 | " + " | ".join(engines) + " |")
        lines.append(f"|-----|------|" + "------|" * len(engines))
        for m in prop["monthly_data"]:
            row = f"| {m['month']} | {m['total_sessions']:,} | "
            row += " | ".join(str(m["by_engine"].get(e, {}).get("sessions", 0)) for e in engines)
            row += " |"
            lines.append(row)
        lines.append(f"")

    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"  保存: {out_path}")


# ── メイン ────────────────────────────────────────────────────
def main():
    print("=== AI参照トラフィック レポート生成 ===\n")

    data = load_data()
    now  = datetime.now().strftime("%Y%m")

    md_path   = os.path.join(BASE_DIR, f"ai_traffic_report_{now}.md")
    pptx_path = os.path.join(BASE_DIR, f"ai_traffic_report_{now}.pptx")

    print("▶ Markdownレポート生成中...", flush=True)
    build_markdown(data, md_path)

    print("▶ PowerPointレポート生成中...", flush=True)
    build_pptx(data, pptx_path)

    print(f"\n✅ レポートを生成しました。")
    print(f"   {md_path}")
    print(f"   {pptx_path}")


if __name__ == "__main__":
    main()
